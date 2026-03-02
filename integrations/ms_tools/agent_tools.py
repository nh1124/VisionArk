"""MS Office integration tools for VisionArk agents.

Tools
-----
RenderPdfTool       Markdown / HTML / text / JSON-template → PDF (via Playwright)
WordTool            Create and edit Word (.docx) documents
ExcelTool           Read / write Excel (.xlsx) workbooks
PptTool             Create and update PowerPoint (.pptx) presentations
MsAuthManagerTool   Microsoft Graph authentication and audit management
"""
from __future__ import annotations

import json
import logging
import os
import re
import time
from pathlib import Path
from typing import Any, Optional

from pydantic import Field
from va_sdk import BaseTool, BaseModel, IntegrationContext, ToolResult

logger = logging.getLogger(__name__)


# ── shared helpers ────────────────────────────────────────────────────

def _artifacts_dir(ctx: IntegrationContext) -> Path:
    from shared.paths import get_project_dir
    if not ctx.project_id:
        raise ValueError("project_id is required for file operations")
    d = get_project_dir(ctx.user_id, ctx.project_id) / "artifacts"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _safe_name(name: str, fallback: str = "output") -> str:
    return re.sub(r'[<>:"/\\|?*]', "_", name).strip() or fallback


def _ok(data: dict) -> ToolResult:
    return ToolResult(content=json.dumps(data), data=data, is_success=True)


def _err(message: str) -> ToolResult:
    return ToolResult(content=message, is_success=False)


def _as_dict(v: Any) -> dict:
    """Coerce a value to dict — handles the case where LLMs stringify JSON objects."""
    if isinstance(v, dict):
        return v
    if isinstance(v, str):
        try:
            result = json.loads(v)
            return result if isinstance(result, dict) else {}
        except (json.JSONDecodeError, ValueError):
            return {}
    return {}


def _as_list(v: Any) -> list:
    """Coerce a value to list — handles the case where LLMs stringify JSON arrays."""
    if isinstance(v, list):
        return v
    if isinstance(v, str):
        try:
            result = json.loads(v)
            return result if isinstance(result, list) else []
        except (json.JSONDecodeError, ValueError):
            return []
    return []


# ═══════════════════════════════════════════════════════════════════════
# render_pdf
# ═══════════════════════════════════════════════════════════════════════

_DEFAULT_CSS = """
body {
  font-family: 'Segoe UI', 'Hiragino Sans', 'Meiryo', sans-serif;
  font-size: 11pt; line-height: 1.7; color: #1a1a1a;
}
h1 { font-size: 20pt; border-bottom: 2px solid #333; padding-bottom: 4px; }
h2 { font-size: 15pt; border-bottom: 1px solid #aaa; padding-bottom: 2px; }
h3 { font-size: 13pt; }
table { border-collapse: collapse; width: 100%; margin: 12px 0; }
th, td { border: 1px solid #ccc; padding: 6px 10px; }
th { background: #f0f0f0; font-weight: bold; }
pre, code { font-family: 'Consolas', monospace; background: #f5f5f5; padding: 2px 4px; }
pre { padding: 10px; }
blockquote { border-left: 3px solid #aaa; margin: 0; padding-left: 12px; color: #555; }
"""

_PLAYWRIGHT_FORMATS = {"A4", "Letter", "Legal", "Tabloid", "A3", "A5"}


class RenderPdfArgs(BaseModel):
    source_type: str = Field(..., description="markdown | html | text | json_template")
    source_content: str = Field(..., description="Content to render as PDF")
    output_filename: str = Field("output", description="Output file name without .pdf extension")
    options: dict = Field(
        default_factory=dict,
        description=(
            "Rendering options: page_size (A4/Letter/A3/A5), "
            "margins ({top,bottom,left,right} in mm), "
            "header (str), footer (str), page_numbers (bool), "
            "toc (bool), extra_css (str)"
        ),
    )
    dry_run: bool = Field(False, description="Validate without writing file")


class RenderPdfTool(BaseTool):
    name = "render_pdf"
    description = (
        "Render content (Markdown, HTML, plain text, or JSON template) to a "
        "production-quality PDF file. Uses Playwright for rendering — supports "
        "full Unicode including Japanese. Saves to artifacts/pdfs/."
    )
    args_schema = RenderPdfArgs

    async def run(self, ctx: IntegrationContext, **kwargs) -> ToolResult:
        source_type: str = kwargs.get("source_type", "text")
        source_content: str = kwargs.get("source_content", "")
        output_filename: str = _safe_name(kwargs.get("output_filename", "output"))
        options: dict = kwargs.get("options", {})
        dry_run: bool = kwargs.get("dry_run", False)

        warnings: list[str] = []

        try:
            html = self._to_html(source_type, source_content, options, warnings)
        except Exception as exc:
            return _err(f"Content conversion failed: {exc}")

        if dry_run:
            return _ok({
                "success": True,
                "dry_run": True,
                "warnings": warnings,
                "message": "Dry run OK — use dry_run=false to write the PDF.",
                "validation_report": self._validate(source_content, warnings),
            })

        try:
            pdf_bytes = await self._render(html, options)
        except ImportError:
            return _err(
                "Playwright/Chromium not available. "
                "Run: pip install playwright && playwright install chromium"
            )
        except Exception as exc:
            logger.exception("render_pdf rendering failed")
            return _err(f"PDF rendering failed: {exc}")

        try:
            art = _artifacts_dir(ctx)
            pdf_dir = art / "pdfs"
            pdf_dir.mkdir(parents=True, exist_ok=True)
            out_path = pdf_dir / f"{output_filename}.pdf"
            out_path.write_bytes(pdf_bytes)
        except Exception as exc:
            return _err(f"Failed to save PDF: {exc}")

        page_count = pdf_bytes.count(b"/Type /Page\n") or pdf_bytes.count(b"/Type/Page")
        if page_count == 0:
            page_count = max(1, len(pdf_bytes) // 50_000)

        return _ok({
            "success": True,
            "pdf_path": f"artifacts/pdfs/{output_filename}.pdf",
            "page_count": page_count,
            "file_size_bytes": len(pdf_bytes),
            "warnings": warnings,
            "validation_report": self._validate(source_content, warnings),
        })

    # ── helpers ──────────────────────────────────────────────────────

    def _to_html(self, source_type: str, content: str, options: dict, warnings: list) -> str:
        if source_type == "markdown":
            body = self._md_to_html(content, warnings)
        elif source_type == "html":
            body = content
        elif source_type == "json_template":
            body = self._json_to_html(content, warnings)
        else:
            escaped = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            body = f"<pre>{escaped}</pre>"

        if options.get("toc"):
            body = self._build_toc(body) + body

        extra_css = options.get("extra_css", "")
        locale = options.get("locale", "")
        lang = f' lang="{locale}"' if locale else ""
        return (
            f'<!DOCTYPE html><html{lang}><head><meta charset="UTF-8">'
            f"<style>{_DEFAULT_CSS}{extra_css}</style></head><body>{body}</body></html>"
        )

    def _md_to_html(self, content: str, warnings: list) -> str:
        try:
            import markdown as md_lib
            return md_lib.markdown(content, extensions=["tables", "fenced_code", "toc", "nl2br"])
        except ImportError:
            warnings.append("markdown library missing; rendering as plain text.")
            return f"<pre>{content}</pre>"
        except Exception as exc:
            warnings.append(f"Markdown warning: {exc}")
            return f"<pre>{content}</pre>"

    def _json_to_html(self, content: str, warnings: list) -> str:
        try:
            data = json.loads(content)
        except json.JSONDecodeError as exc:
            warnings.append(f"Invalid JSON: {exc}")
            return f"<pre>{content}</pre>"
        parts: list[str] = []
        if t := data.get("title"):
            parts.append(f"<h1>{t}</h1>")
        if s := data.get("subtitle"):
            parts.append(f"<p><em>{s}</em></p>")
        for section in data.get("sections", []):
            if h := section.get("heading"):
                parts.append(f"<h2>{h}</h2>")
            if b := section.get("body"):
                parts.append(f"<p>{b}</p>")
            if items := section.get("items", []):
                parts.append("<ul>" + "".join(f"<li>{i}</li>" for i in items) + "</ul>")
            headers = section.get("table_headers", [])
            rows = section.get("table_rows", [])
            if headers or rows:
                parts.append("<table>")
                if headers:
                    parts.append("<thead><tr>" + "".join(f"<th>{h}</th>" for h in headers) + "</tr></thead>")
                if rows:
                    parts.append("<tbody>" + "".join(
                        "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>" for row in rows
                    ) + "</tbody>")
                parts.append("</table>")
        return "\n".join(parts) or "<p>(empty)</p>"

    def _build_toc(self, html_body: str) -> str:
        headings = re.findall(r"<h([1-3])[^>]*?>(.*?)</h\1>", html_body, re.I | re.S)
        if not headings:
            return ""
        items = [
            f"<li>{'&nbsp;' * (int(lvl) - 1) * 4}{re.sub(r'<[^>]+>', '', txt).strip()}</li>"
            for lvl, txt in headings
        ]
        return f"<h2>Table of Contents</h2><ul>{''.join(items)}</ul><hr>"

    async def _render(self, html: str, options: dict) -> bytes:
        from playwright.async_api import async_playwright
        page_size = options.get("page_size", "A4")
        if page_size not in _PLAYWRIGHT_FORMATS:
            page_size = "A4"
        m = options.get("margins", {})
        header_text = options.get("header", "")
        footer_text = options.get("footer", "")
        page_numbers = options.get("page_numbers", False)
        _s = "font-family:'Segoe UI',sans-serif;font-size:9px;color:#555;width:100%;padding:0 20mm;box-sizing:border-box;"
        display_hf = bool(header_text or footer_text or page_numbers)
        header_tmpl = f'<div style="{_s}text-align:center;">{header_text}</div>' if header_text else ("<span></span>" if display_hf else "")
        pn = '<span class="pageNumber"></span> / <span class="totalPages"></span>' if page_numbers else ""
        footer_body = "  ·  ".join(filter(None, [footer_text, pn]))
        footer_tmpl = f'<div style="{_s}text-align:center;">{footer_body}</div>' if footer_body else ("<span></span>" if display_hf else "")

        async with async_playwright() as pw:
            browser = await pw.chromium.launch(args=["--no-sandbox"])
            try:
                page = await browser.new_page()
                await page.set_content(html, wait_until="networkidle")
                pdf_bytes = await page.pdf(
                    format=page_size,
                    margin={
                        "top": f"{m.get('top', 20)}mm",
                        "bottom": f"{m.get('bottom', 20)}mm",
                        "left": f"{m.get('left', 20)}mm",
                        "right": f"{m.get('right', 20)}mm",
                    },
                    display_header_footer=display_hf,
                    header_template=header_tmpl,
                    footer_template=footer_tmpl,
                    print_background=True,
                )
            finally:
                await browser.close()
        return pdf_bytes

    def _validate(self, content: str, warnings: list) -> dict:
        return {
            "missing_assets": re.findall(r'!\[.*?\]\(((?!https?://).+?)\)', content),
            "unresolved_links": re.findall(r'\[.*?\]\(((?!https?://).+?)\)', content),
            "warnings_count": len(warnings),
        }


# ═══════════════════════════════════════════════════════════════════════
# word_tool
# ═══════════════════════════════════════════════════════════════════════

class WordArgs(BaseModel):
    operation: str = Field(
        ...,
        description=(
            "Operation to perform. Valid values: "
            "'create_from_template' (alias: 'create') — create a new .docx; "
            "'patch_document' (alias: 'patch') — apply ops to existing doc; "
            "'apply_styles' — map text to Word style names; "
            "'add_comments' — append review comments; "
            "'export' — save as docx or pdf."
        ),
    )
    doc_id: Optional[str] = Field(None, description="Relative path to existing .docx under artifacts/")
    template_id: Optional[str] = Field("blank", description="'blank','report','memo' or path to .docx template")
    output_filename: Optional[str] = Field("document", description="Output filename without extension")
    data: Optional[dict] = Field(None, description="Content dict for create: {title, sections:[{heading,paragraphs[],level}], tables:[{headers[],rows[][]}]}")
    operations: Optional[list] = Field(None, description="Patch ops: [{op:'replace'|'insert'|'delete', target, value}]")
    style_map: Optional[dict] = Field(None, description="Text → style name mapping for apply_styles")
    comments: Optional[list] = Field(None, description="Comment specs: [{text, anchor, author}]")
    export_format: Optional[str] = Field("docx", description="'docx' or 'pdf'")


class WordTool(BaseTool):
    name = "word_tool"
    description = (
        "Create and edit Microsoft Word (.docx) documents. "
        "Supports template-based creation, section patching, style application, "
        "review comments, and export to .docx or PDF. "
        "Files saved under artifacts/documents/."
    )
    args_schema = WordArgs

    _OP_ALIASES: dict[str, str] = {
        "create": "create_from_template",
        "patch": "patch_document",
        "update": "patch_document",
        "style": "apply_styles",
        "comment": "add_comments",
    }

    async def run(self, ctx: IntegrationContext, **kwargs) -> ToolResult:
        op = self._OP_ALIASES.get(kwargs.get("operation", ""), kwargs.get("operation", ""))
        try:
            import docx  # noqa
        except ImportError:
            return _err("python-docx not installed. Run: pip install python-docx")

        try:
            if op == "create_from_template":
                return await self._create(ctx, kwargs)
            if op == "patch_document":
                return await self._patch(ctx, kwargs)
            if op == "apply_styles":
                return await self._apply_styles(ctx, kwargs)
            if op == "add_comments":
                return await self._add_comments(ctx, kwargs)
            if op == "export":
                return await self._export(ctx, kwargs)
            return _err(
                f"Unknown operation: {op!r}. "
                f"Valid values: create_from_template (or 'create'), patch_document (or 'patch'), "
                f"apply_styles, add_comments, export."
            )
        except Exception as exc:
            logger.exception("word_tool.%s failed", op)
            return _err(f"word_tool.{op} failed: {exc}")

    async def _create(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from docx import Document
        data: dict = _as_dict(kw.get("data") or {})
        filename = _safe_name(kw.get("output_filename") or "document", "document")
        art = _artifacts_dir(ctx)
        doc_dir = art / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)

        template_path = _resolve_template(kw.get("template_id") or "blank", art)
        doc = Document(str(template_path)) if template_path else Document()

        if title := data.get("title"):
            doc.add_heading(title, level=0)
        for section in data.get("sections", []):
            if h := section.get("heading"):
                doc.add_heading(h, level=section.get("level", 1))
            for para in section.get("paragraphs", []):
                doc.add_paragraph(para)
        for tbl in data.get("tables", []):
            headers = tbl.get("headers", [])
            rows = tbl.get("rows", [])
            cols = max(len(headers), max((len(r) for r in rows), default=0))
            if not cols:
                continue
            t = doc.add_table(rows=1 + len(rows), cols=cols)
            t.style = "Table Grid"
            for i, h in enumerate(headers[:cols]):
                t.cell(0, i).text = str(h)
            for ri, row in enumerate(rows):
                for ci, cell in enumerate(row[:cols]):
                    t.cell(ri + 1, ci).text = str(cell)

        out = doc_dir / f"{filename}.docx"
        doc.save(str(out))
        return _ok({"success": True, "doc_id": f"artifacts/documents/{filename}.docx",
                    "paragraphs": len(doc.paragraphs)})

    async def _patch(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from docx import Document
        art = _artifacts_dir(ctx)
        doc_path = _resolve_doc(kw.get("doc_id") or "", art)
        if not doc_path:
            return _err(f"Document not found: {kw.get('doc_id')!r}")
        doc = Document(str(doc_path))
        log: list[str] = []
        for op_spec in _as_list(kw.get("operations") or []):
            op, target, value = op_spec.get("op"), op_spec.get("target", ""), op_spec.get("value", "")
            if op == "replace":
                n = sum(1 for p in doc.paragraphs if target in p.text
                        for r in p.runs if target in r.text and not setattr(r, "text", r.text.replace(target, value)))
                log.append(f"replace: {n} run(s)")
            elif op == "insert":
                doc.add_paragraph(value)
                log.append("insert: paragraph appended")
            elif op == "delete":
                removed = 0
                for p in list(doc.paragraphs):
                    if target in p.text:
                        p._element.getparent().remove(p._element)
                        removed += 1
                log.append(f"delete: {removed} paragraph(s)")
        doc.save(str(doc_path))
        return _ok({"success": True, "doc_id": kw.get("doc_id"), "changes": log})

    async def _apply_styles(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from docx import Document
        art = _artifacts_dir(ctx)
        doc_path = _resolve_doc(kw.get("doc_id") or "", art)
        if not doc_path:
            return _err(f"Document not found: {kw.get('doc_id')!r}")
        style_map: dict = _as_dict(kw.get("style_map") or {})
        doc = Document(str(doc_path))
        applied: list[str] = []
        for para in doc.paragraphs:
            if (sname := style_map.get(para.text.strip())):
                try:
                    para.style = doc.styles[sname]
                    applied.append(f"'{para.text.strip()}' → {sname}")
                except KeyError:
                    applied.append(f"WARN: style '{sname}' not found")
        doc.save(str(doc_path))
        return _ok({"success": True, "doc_id": kw.get("doc_id"), "styles_applied": len(applied), "log": applied})

    async def _add_comments(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from docx import Document
        art = _artifacts_dir(ctx)
        doc_path = _resolve_doc(kw.get("doc_id") or "", art)
        if not doc_path:
            return _err(f"Document not found: {kw.get('doc_id')!r}")
        doc = Document(str(doc_path))
        doc.add_paragraph("")
        doc.add_heading("Review Comments", level=2)
        added = 0
        for c in _as_list(kw.get("comments") or []):
            doc.add_paragraph(f"[{c.get('author','Reviewer')}] Re: '{c.get('anchor','')}' — {c.get('text','')}")
            added += 1
        doc.save(str(doc_path))
        return _ok({"success": True, "doc_id": kw.get("doc_id"), "comments_added": added,
                    "note": "Comments appended as a 'Review Comments' section."})

    async def _export(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        fmt = kw.get("export_format", "docx")
        doc_id = kw.get("doc_id") or ""
        if fmt == "docx":
            return _ok({"success": True, "export_path": doc_id, "format": "docx",
                        "message": "File is already in .docx format."})
        if fmt == "pdf":
            art = _artifacts_dir(ctx)
            doc_path = _resolve_doc(doc_id, art)
            if not doc_path:
                return _err(f"Document not found: {doc_id!r}")
            from docx import Document
            doc = Document(str(doc_path))
            text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            pdf_tool = RenderPdfTool()
            return await pdf_tool.run(
                ctx,
                source_type="text",
                source_content=text,
                output_filename=doc_path.stem,
                options={"page_size": "A4"},
            )
        return _err(f"Unsupported export format: {fmt!r}")


# ═══════════════════════════════════════════════════════════════════════
# excel_tool
# ═══════════════════════════════════════════════════════════════════════

class ExcelArgs(BaseModel):
    operation: str = Field(
        ...,
        description=(
            "Operation to perform. Valid values: "
            "'read_range' (alias: 'read') — read cells; "
            "'update_range' (alias: 'write' or 'create') — write cell values; "
            "'apply_formula' (alias: 'formula') — set a formula in a range; "
            "'create_chart' (alias: 'chart') — add a chart to a sheet; "
            "'validate' — check data against rules."
        ),
    )
    workbook_id: Optional[str] = Field("workbook.xlsx", description="Path to .xlsx under artifacts/ or new filename")
    sheet: Optional[str] = Field(None, description="Worksheet name (default: active sheet)")
    range: Optional[str] = Field(None, description="A1-notation range, e.g. 'A1:D10'")
    values: Optional[list] = Field(None, description="2-D array of values for update_range")
    formula: Optional[str] = Field(None, description="Formula string for apply_formula, e.g. '=SUM(B2:B10)'")
    chart_config: Optional[dict] = Field(None, description="Chart config: {type,data_range,title,x_title,y_title,position}")
    validation_rules: Optional[list] = Field(None, description="Validation rules: [{range,type,expected_type,min,max,pattern}]")


class ExcelTool(BaseTool):
    name = "excel_tool"
    description = (
        "Read and write Microsoft Excel (.xlsx) workbooks. "
        "Supports range reads/writes, formulas, bar/line/pie charts, and data validation. "
        "Files saved under artifacts/spreadsheets/."
    )
    args_schema = ExcelArgs

    _OP_ALIASES: dict[str, str] = {
        "read": "read_range",
        "write": "update_range",
        "create": "update_range",
        "update": "update_range",
        "formula": "apply_formula",
        "chart": "create_chart",
        "add_chart": "create_chart",
    }

    async def run(self, ctx: IntegrationContext, **kwargs) -> ToolResult:
        op = self._OP_ALIASES.get(kwargs.get("operation", ""), kwargs.get("operation", ""))
        try:
            import openpyxl  # noqa
        except ImportError:
            return _err("openpyxl not installed. Run: pip install openpyxl")
        try:
            if op == "read_range":
                return await self._read(ctx, kwargs)
            if op == "update_range":
                return await self._update(ctx, kwargs)
            if op == "apply_formula":
                return await self._formula(ctx, kwargs)
            if op == "create_chart":
                return await self._chart(ctx, kwargs)
            if op == "validate":
                return await self._validate(ctx, kwargs)
            return _err(
                f"Unknown operation: {op!r}. "
                f"Valid values: read_range (or 'read'), update_range (or 'write'/'create'), "
                f"apply_formula (or 'formula'), create_chart (or 'chart'), validate."
            )
        except Exception as exc:
            logger.exception("excel_tool.%s failed", op)
            return _err(f"excel_tool.{op} failed: {exc}")

    async def _read(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        import openpyxl
        art = _artifacts_dir(ctx)
        wb_path = _resolve_wb(kw.get("workbook_id") or "workbook.xlsx", art)
        if not wb_path or not wb_path.exists():
            return _err(f"Workbook not found: {kw.get('workbook_id')!r}")
        wb = openpyxl.load_workbook(str(wb_path), data_only=True)
        ws = _get_ws(wb, kw.get("sheet"))
        rng = kw.get("range")
        data = _cells_to_list(ws[rng]) if rng else [[c.value for c in row] for row in ws.iter_rows()]
        return _ok({"success": True, "sheet": ws.title, "range": rng or "all",
                    "rows": len(data), "values": data})

    async def _update(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        import openpyxl
        values: list = _as_list(kw.get("values") or [])
        if not values:
            return _err("values must be a non-empty 2-D array.")
        art = _artifacts_dir(ctx)
        ss_dir = art / "spreadsheets"
        ss_dir.mkdir(parents=True, exist_ok=True)
        wb_id = kw.get("workbook_id") or "workbook.xlsx"
        wb_path = _resolve_wb(wb_id, art)
        wb = openpyxl.load_workbook(str(wb_path)) if (wb_path and wb_path.exists()) else openpyxl.Workbook()
        if not (wb_path and wb_path.exists()):
            wb_path = ss_dir / _safe_wb_name(wb_id)
        ws = _get_ws(wb, kw.get("sheet"))
        start = (kw.get("range") or "A1").split(":")[0]
        col_m = re.match(r"([A-Za-z]+)", start)
        row_m = re.search(r"(\d+)", start)
        if not col_m or not row_m:
            return _err(f"Invalid range: {kw.get('range')!r}")
        sc = openpyxl.utils.column_index_from_string(col_m.group(1))
        sr = int(row_m.group(1))
        written = sum(
            1 for ri, row in enumerate(values)
            for ci, val in enumerate(row)
            if ws.cell(row=sr + ri, column=sc + ci, value=val)
        )
        wb.save(str(wb_path))
        rel = str(wb_path.relative_to(art.parent))
        return _ok({"success": True, "workbook_id": rel, "sheet": ws.title, "cells_written": written})

    async def _formula(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        import openpyxl
        formula = kw.get("formula") or ""
        if not formula:
            return _err("formula is required.")
        art = _artifacts_dir(ctx)
        wb_path = _resolve_wb(kw.get("workbook_id") or "workbook.xlsx", art)
        if not wb_path or not wb_path.exists():
            return _err(f"Workbook not found: {kw.get('workbook_id')!r}")
        wb = openpyxl.load_workbook(str(wb_path))
        ws = _get_ws(wb, kw.get("sheet"))
        rng = kw.get("range") or "A1"
        applied: list[str] = []
        for row in ws[rng]:
            cells = row if hasattr(row, "__iter__") and not hasattr(row, "value") else [row]
            for cell in cells:
                cell.value = formula
                applied.append(cell.coordinate)
        wb.save(str(wb_path))
        return _ok({"success": True, "formula": formula, "applied_to": applied})

    async def _chart(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        import openpyxl
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        cfg: dict = _as_dict(kw.get("chart_config") or {})
        art = _artifacts_dir(ctx)
        wb_path = _resolve_wb(kw.get("workbook_id") or "workbook.xlsx", art)
        if not wb_path or not wb_path.exists():
            return _err(f"Workbook not found: {kw.get('workbook_id')!r}")
        wb = openpyxl.load_workbook(str(wb_path))
        ws = _get_ws(wb, kw.get("sheet"))
        ctype = cfg.get("type", "bar").lower()
        chart = {"line": LineChart, "pie": PieChart}.get(ctype, BarChart)()
        chart.title = cfg.get("title", "Chart")
        if xt := cfg.get("x_title"):
            chart.x_axis.title = xt
        if yt := cfg.get("y_title"):
            chart.y_axis.title = yt
        try:
            ref = Reference(ws, range_string=f"{ws.title}!{cfg.get('data_range','A1:B5')}")
        except Exception as exc:
            return _err(f"Invalid data_range: {exc}")
        chart.add_data(ref, titles_from_data=True)
        ws.add_chart(chart, cfg.get("position", "E2"))
        wb.save(str(wb_path))
        return _ok({"success": True, "chart_type": ctype, "title": chart.title,
                    "position": cfg.get("position", "E2")})

    async def _validate(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        import openpyxl
        art = _artifacts_dir(ctx)
        wb_path = _resolve_wb(kw.get("workbook_id") or "workbook.xlsx", art)
        if not wb_path or not wb_path.exists():
            return _err(f"Workbook not found: {kw.get('workbook_id')!r}")
        wb = openpyxl.load_workbook(str(wb_path), data_only=True)
        ws = _get_ws(wb, kw.get("sheet"))
        violations: list[dict] = []
        for rule in _as_list(kw.get("validation_rules") or []):
            rng, rtype = rule.get("range", "A1"), rule.get("type", "type")
            for row in ws[rng]:
                cells = row if hasattr(row, "__iter__") and not hasattr(row, "value") else [row]
                for cell in cells:
                    val, coord = cell.value, cell.coordinate
                    if rtype == "type":
                        exp = rule.get("expected_type", "")
                        if exp == "number" and not isinstance(val, (int, float)):
                            violations.append({"cell": coord, "issue": f"Expected number, got {type(val).__name__}"})
                        elif exp == "string" and not isinstance(val, str):
                            violations.append({"cell": coord, "issue": f"Expected string, got {type(val).__name__}"})
                    elif rtype == "range":
                        if isinstance(val, (int, float)):
                            mn, mx = rule.get("min"), rule.get("max")
                            if mn is not None and val < mn:
                                violations.append({"cell": coord, "issue": f"{val} < min {mn}"})
                            if mx is not None and val > mx:
                                violations.append({"cell": coord, "issue": f"{val} > max {mx}"})
                    elif rtype == "regex":
                        pat = rule.get("pattern", "")
                        if val is not None and not re.match(pat, str(val)):
                            violations.append({"cell": coord, "issue": f"'{val}' !~ /{pat}/"})
        return _ok({"success": True, "valid": not violations, "violations": violations,
                    "violation_count": len(violations),
                    "recovery_hint": "All cells passed." if not violations
                    else "Review listed cells before writing."})


# ═══════════════════════════════════════════════════════════════════════
# ppt_tool
# ═══════════════════════════════════════════════════════════════════════

class PptArgs(BaseModel):
    operation: str = Field(
        ...,
        description=(
            "Operation to perform. Valid values: "
            "'create_from_template' (alias: 'create') — create a new .pptx from an outline; "
            "'upsert_slide' (alias: 'add_slide') — add or update a single slide; "
            "'bind_chart' — embed an Excel chart into a slide; "
            "'export' — export to pptx or pdf."
        ),
    )
    deck_id: Optional[str] = Field(None, description="Relative path to existing .pptx under artifacts/")
    template_id: Optional[str] = Field("blank", description="'blank','title_content','two_column' or path to template")
    output_filename: Optional[str] = Field(
        "presentation",
        description="Output filename without extension (alias: file_path — strip .pptx suffix if provided)",
    )
    outline: Optional[list] = Field(
        None,
        description=(
            "Slide specs for create_from_template (alias: 'slides'). "
            "Each item: {title, bullets: [], notes: '', layout: ''}"
        ),
    )
    slide_spec: Optional[dict] = Field(None, description="Slide spec for upsert: {index,title,bullets[],layout,notes}")
    excel_source: Optional[dict] = Field(None, description="Excel data for bind_chart: {workbook_id,sheet,data_range,chart_type,title,slide_index}")
    export_format: Optional[str] = Field("pptx", description="'pptx' or 'pdf'")


class PptTool(BaseTool):
    name = "ppt_tool"
    description = (
        "Create and update Microsoft PowerPoint (.pptx) presentations. "
        "Use operation='create_from_template' (or 'create') with outline=[{title,bullets}] to build a new deck. "
        "Use operation='upsert_slide' (or 'add_slide') with slide_spec to add/edit one slide. "
        "Use operation='bind_chart' to embed an Excel chart. "
        "Use operation='export' to save as pptx or pdf. "
        "Files saved under artifacts/presentations/."
    )
    args_schema = PptArgs

    # Aliases accepted from LLMs that use shorter operation names
    _OP_ALIASES: dict[str, str] = {
        "create": "create_from_template",
        "add_slide": "upsert_slide",
        "update_slide": "upsert_slide",
        "chart": "bind_chart",
    }

    async def run(self, ctx: IntegrationContext, **kwargs) -> ToolResult:
        op = kwargs.get("operation", "")
        # Resolve aliases so LLMs can use natural names like "create"
        op = self._OP_ALIASES.get(op, op)

        try:
            import pptx  # noqa
        except ImportError:
            return _err("python-pptx not installed. Run: pip install python-pptx")
        try:
            if op == "create_from_template":
                return await self._create(ctx, kwargs)
            if op == "upsert_slide":
                return await self._upsert(ctx, kwargs)
            if op == "bind_chart":
                return await self._bind_chart(ctx, kwargs)
            if op == "export":
                return await self._export(ctx, kwargs)
            return _err(
                f"Unknown operation: {op!r}. "
                f"Valid values: create_from_template (or 'create'), upsert_slide (or 'add_slide'), bind_chart, export."
            )
        except Exception as exc:
            logger.exception("ppt_tool.%s failed", op)
            return _err(f"ppt_tool.{op} failed: {exc}")

    async def _create(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from pptx import Presentation
        art = _artifacts_dir(ctx)
        pres_dir = art / "presentations"
        pres_dir.mkdir(parents=True, exist_ok=True)

        # Accept both 'output_filename' and 'file_path' (strip extension if provided)
        raw_name = kw.get("output_filename") or kw.get("file_path") or "presentation"
        raw_name = Path(raw_name).stem  # strip .pptx / .ppt suffix if given
        filename = _safe_name(raw_name, "presentation")

        tmpl = _resolve_pptx_template(kw.get("template_id") or "blank", art)
        prs = Presentation(str(tmpl)) if tmpl else Presentation()

        # Accept both 'outline' and 'slides' as the slide list key
        slides_raw = kw.get("outline") or kw.get("slides") or []
        for spec in _as_list(slides_raw):
            _add_slide(prs, spec)

        out = pres_dir / f"{filename}.pptx"
        prs.save(str(out))
        return _ok({"success": True, "deck_id": f"artifacts/presentations/{filename}.pptx",
                    "slide_count": len(prs.slides)})

    async def _upsert(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from pptx import Presentation
        art = _artifacts_dir(ctx)
        deck_path = _resolve_deck(kw.get("deck_id") or "", art)
        if not deck_path:
            return _err(f"Presentation not found: {kw.get('deck_id')!r}")
        prs = Presentation(str(deck_path))
        spec: dict = _as_dict(kw.get("slide_spec") or {})
        idx = spec.get("index")
        if idx is not None and 0 <= idx < len(prs.slides):
            _update_slide(prs.slides[idx], spec)
            action = "updated"
        else:
            _add_slide(prs, spec)
            idx = len(prs.slides) - 1
            action = "appended"
        prs.save(str(deck_path))
        return _ok({"success": True, "deck_id": kw.get("deck_id"), "action": action,
                    "slide_index": idx, "slide_count": len(prs.slides)})

    async def _bind_chart(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        import openpyxl

        src: dict = _as_dict(kw.get("excel_source") or {})
        art = _artifacts_dir(ctx)
        deck_path = _resolve_deck(kw.get("deck_id") or "", art)
        if not deck_path:
            return _err(f"Presentation not found: {kw.get('deck_id')!r}")

        wb_path = _resolve_wb(src.get("workbook_id") or "", art)
        if not wb_path or not wb_path.exists():
            return _err(f"Workbook not found: {src.get('workbook_id')!r}")

        wb = openpyxl.load_workbook(str(wb_path), data_only=True)
        ws = _get_ws(wb, src.get("sheet"))
        raw = _cells_to_list(ws[src.get("data_range", "A1:B5")])
        if len(raw) < 2:
            return _err("data_range must have ≥ 2 rows (headers + data)")

        chart_data = ChartData()
        chart_data.categories = [str(r[0]) for r in raw[1:]]
        headers = raw[0]
        for ci in range(1, len(headers) if headers else 1):
            chart_data.add_series(
                str(headers[ci]) if headers[ci] is not None else f"Series {ci}",
                [r[ci] if ci < len(r) else 0 for r in raw[1:]],
            )

        type_map = {"line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}
        xl_type = type_map.get(src.get("chart_type", "bar").lower(), XL_CHART_TYPE.BAR_CLUSTERED)
        prs = Presentation(str(deck_path))
        si = src.get("slide_index", 0)
        if si >= len(prs.slides):
            return _err(f"slide_index {si} out of range")
        chart = prs.slides[si].shapes.add_chart(
            xl_type, Inches(1), Inches(2), Inches(8), Inches(4.5), chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = src.get("title", "Chart")
        prs.save(str(deck_path))
        return _ok({"success": True, "deck_id": kw.get("deck_id"),
                    "chart_type": src.get("chart_type", "bar"), "slide_index": si})

    async def _export(self, ctx: IntegrationContext, kw: dict) -> ToolResult:
        fmt = kw.get("export_format", "pptx")
        deck_id = kw.get("deck_id") or ""
        if fmt == "pptx":
            return _ok({"success": True, "export_path": deck_id, "format": "pptx",
                        "message": "File is already in .pptx format."})
        if fmt == "pdf":
            from pptx import Presentation
            art = _artifacts_dir(ctx)
            deck_path = _resolve_deck(deck_id, art)
            if not deck_path:
                return _err(f"Presentation not found: {deck_id!r}")
            prs = Presentation(str(deck_path))
            sections = []
            for i, slide in enumerate(prs.slides):
                texts = [sh.text_frame.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
                sections.append({"heading": texts[0] if texts else f"Slide {i+1}", "paragraphs": texts[1:]})
            content = json.dumps({"title": deck_path.stem, "sections": sections})
            return await RenderPdfTool().run(ctx, source_type="json_template",
                                             source_content=content,
                                             output_filename=deck_path.stem,
                                             options={"page_size": "A4"})
        return _err(f"Unsupported export format: {fmt!r}")


# ═══════════════════════════════════════════════════════════════════════
# ms_auth_manager
# ═══════════════════════════════════════════════════════════════════════

class MsAuthArgs(BaseModel):
    operation: str = Field(..., description="get_token | validate_scope | refresh_token | attach_audit_context")
    scopes: Optional[list] = Field(None, description="MS Graph permission scopes, e.g. ['Files.ReadWrite']")
    required_scopes: Optional[list] = Field(None, description="Scopes to verify for validate_scope")
    operation_id: Optional[str] = Field(None, description="Operation ID for audit context")
    actor: Optional[str] = Field(None, description="Actor name for audit entry")
    resource: Optional[str] = Field(None, description="Resource URI being accessed")


class MsAuthManagerTool(BaseTool):
    name = "ms_auth_manager"
    description = (
        "Manage Microsoft Graph authentication and authorization. "
        "Acquire tokens (client credentials or user-delegated), validate scopes, "
        "refresh expired tokens, and attach audit context to operations. "
        "Configure via AZURE_CLIENT_ID / AZURE_CLIENT_SECRET / AZURE_TENANT_ID."
    )
    args_schema = MsAuthArgs

    async def run(self, ctx: IntegrationContext, **kwargs) -> ToolResult:
        from integrations.ms_tools._auth import get_token_from_ctx, clear_token_cache
        op = kwargs.get("operation", "")

        if op == "get_token":
            scopes = kwargs.get("scopes") or ["https://graph.microsoft.com/.default"]
            injected = (ctx.metadata or {}).get("ms_access_token")
            if injected:
                return _ok({"success": True, "token_source": "user_delegated",
                            "token_preview": f"{injected[:12]}…",
                            "message": "Using user-delegated token from context."})
            token = get_token_from_ctx(ctx, scopes)
            if token:
                return _ok({"success": True, "token_source": "client_credentials",
                            "token_preview": f"{token[:12]}…",
                            "message": "Token acquired via client credentials."})
            return _ok({"success": False, "token_source": "none",
                        "message": "No MS credentials configured. Set AZURE_CLIENT_ID/SECRET/TENANT_ID. "
                                   "Office tools operate in offline/local-file mode."})

        if op == "validate_scope":
            required = kwargs.get("required_scopes") or []
            if not required:
                return _ok({"success": True, "valid": True, "message": "No scopes to validate."})
            token = get_token_from_ctx(ctx, ["https://graph.microsoft.com/.default"])
            if not token:
                return _ok({"success": True, "valid": False, "missing_scopes": required,
                            "message": "No token available."})
            granted = self._extract_scopes(token)
            missing = [s for s in required if s not in granted]
            return _ok({"success": True, "valid": not missing, "granted_scopes": granted,
                        "missing_scopes": missing})

        if op == "refresh_token":
            clear_token_cache()
            return _ok({"success": True,
                        "message": "Token cache cleared. Next call will re-acquire."})

        if op == "attach_audit_context":
            entry = {
                "operation_id": kwargs.get("operation_id", ""),
                "actor": kwargs.get("actor", ctx.user_id),
                "resource": kwargs.get("resource", ""),
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            }
            if ctx.metadata is None:
                ctx.metadata = {}
            ctx.metadata.setdefault("ms_audit_log", []).append(entry)
            return _ok({"success": True, "audit_entry": entry,
                        "message": "Audit context attached."})

        return _err(f"Unknown operation: {op!r}")

    def _extract_scopes(self, token: str) -> list[str]:
        try:
            import base64
            parts = token.split(".")
            if len(parts) < 2:
                return []
            payload = json.loads(base64.urlsafe_b64decode(parts[1] + "==").decode("utf-8", errors="replace"))
            scp = payload.get("scp", "")
            roles = payload.get("roles", [])
            return (scp.split() if isinstance(scp, str) else []) + (roles if isinstance(roles, list) else [])
        except Exception:
            return []


# ═══════════════════════════════════════════════════════════════════════
# Shared file-resolution helpers
# ═══════════════════════════════════════════════════════════════════════

def _resolve_template(template_id: str, art: Path) -> Path | None:
    if not template_id or template_id in ("blank",):
        return None
    p = art.parent / template_id
    return p if p.exists() else None


def _resolve_doc(doc_id: str, art: Path) -> Path | None:
    if not doc_id:
        return None
    for base in (art.parent, art, art / "documents"):
        p = base / doc_id if not Path(doc_id).is_absolute() else Path(doc_id)
        if p.exists():
            return p
    return None


def _resolve_wb(wb_id: str, art: Path) -> Path | None:
    if not wb_id:
        return None
    for base in (art.parent, art / "spreadsheets"):
        p = base / wb_id
        if p.exists():
            return p
    return art / "spreadsheets" / _safe_wb_name(wb_id)


def _safe_wb_name(name: str) -> str:
    n = re.sub(r'[<>:"/\\|?*]', "_", name).strip() or "workbook"
    return n if n.endswith(".xlsx") else n + ".xlsx"


def _resolve_deck(deck_id: str, art: Path) -> Path | None:
    if not deck_id:
        return None
    for base in (art.parent, art / "presentations"):
        p = base / deck_id
        if p.exists():
            return p
    return None


def _resolve_pptx_template(template_id: str, art: Path) -> Path | None:
    if not template_id or template_id in ("blank", "title_content", "two_column"):
        return None
    p = art.parent / template_id
    return p if p.exists() else None


def _get_ws(wb, sheet_name: str | None):
    if sheet_name and sheet_name in wb.sheetnames:
        return wb[sheet_name]
    return wb.active


def _cells_to_list(cells) -> list[list]:
    from openpyxl.cell.cell import Cell
    if isinstance(cells, Cell):
        return [[cells.value]]
    if cells and isinstance(cells[0], Cell):
        return [[c.value for c in cells]]
    return [[c.value for c in row] for row in cells]


# ── PowerPoint slide helpers ──────────────────────────────────────────

def _pick_layout(prs, layout_name: str):
    """Find a slide layout by name; fall back to a sensible default."""
    normalized = layout_name.lower().replace(" ", "_")
    for layout in prs.slide_layouts:
        if layout.name.lower().replace(" ", "_") == normalized:
            return layout
    fallback_idx = {"title": 0, "content": 1, "title_content": 1, "two_column": 3, "blank": 6}
    idx = fallback_idx.get(layout_name.lower(), 1)
    try:
        return prs.slide_layouts[idx]
    except IndexError:
        return prs.slide_layouts[0]


def _add_slide(prs, slide_spec: dict) -> None:
    """Append a new slide to the presentation from a spec dict."""
    layout = _pick_layout(prs, slide_spec.get("layout", "content"))
    slide = prs.slides.add_slide(layout)

    title_text = slide_spec.get("title", "")
    bullets = _as_list(slide_spec.get("bullets") or [])
    notes_text = slide_spec.get("notes", "")

    if slide.shapes.title and title_text:
        slide.shapes.title.text = title_text

    body_ph = next(
        (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
        None,
    )
    if body_ph and bullets:
        tf = body_ph.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                tf.add_paragraph().text = bullet

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def _update_slide(slide, slide_spec: dict) -> None:
    """Update title and bullets on an existing slide in-place."""
    title_text = slide_spec.get("title")
    bullets = _as_list(slide_spec.get("bullets") or [])

    if title_text is not None and slide.shapes.title:
        slide.shapes.title.text = title_text

    if bullets:
        body_ph = next(
            (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
            None,
        )
        if body_ph:
            tf = body_ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    tf.add_paragraph().text = bullet
